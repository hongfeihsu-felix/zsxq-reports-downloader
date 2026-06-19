#!/usr/bin/env python3
"""
PDF Report Analyzer - 从 PDF 提取文本（OCR fallback），用 LLM 分析投行报告
替代原 vision_analyze (Xiaomi MIMO LLM)

流程：PDF → 尝试直接文本提取 → 失败则 OCR → LLM 分析 → Markdown → vision_parser 结构化

环境变量（复用当前配置）：
  ANTHROPIC_AUTH_TOKEN  - API Key
  ANTHROPIC_BASE_URL    - API 端点

用法：
  python3 pdf_vision_analyzer.py report.pdf
  python3 pdf_vision_analyzer.py report.pdf --scan-only
  python3 pdf_vision_analyzer.py report.pdf --no-scan
  python3 pdf_vision_analyzer.py ./reports/ --batch
"""

import os
import sys
import re
import json
import argparse
from pathlib import Path
from datetime import datetime

import fitz  # PyMuPDF
import anthropic
import pytesseract
from PIL import Image
import io


# ============ 配置 ============
API_KEY = os.environ.get("ANTHROPIC_AUTH_TOKEN", "")
BASE_URL = os.environ.get("ANTHROPIC_BASE_URL", "https://api.deepseek.com/anthropic")
MODEL = os.environ.get("ANTHROPIC_MODEL", "deepseek-v4-pro[1m]")

SCORE_THRESHOLD = 10
SCANNER_MAX_CHARS = 6000
DEEP_MAX_CHARS = 60000
OCR_DPI = 200  # OCR 渲染 DPI

# ============ 系统提示词 ============

SYSTEM_PROMPT_SCANNER = """You are a financial analyst screening investment banking research reports for a semiconductor/tech hardware fund.

Analyze the extracted text from the first pages of this report and output exactly this format:

## Scanner Report
- **Company**: [full company name, or "Unknown"]
- **Ticker**: [stock ticker, or "N/A"]
- **Rating**: [Buy/Neutral/Sell, or "N/A"]
- **Topic**: [1-sentence summary of what this report covers]
- **Relevance Score**: [0-20]

Scoring rubric (semiconductor/tech hardware focus):
- +5: Covers semiconductor companies (TSMC, NVIDIA, AMD, Intel, Broadcom, Qualcomm, MediaTek, Marvell, Micron, SK Hynix, Samsung, etc.)
- +5: Discusses AI chips, data center GPU/TPU/ASIC, HBM, foundry, advanced packaging
- +3: Includes target price or estimate revisions
- +3: Analyzes industry supply/demand, capex, utilization rates
- +2: Covers optical modules, interconnects, power semiconductors
- +2: Provides revenue/profit forecasts with specific numbers"""

SYSTEM_PROMPT_DEEP = """你是一位顶级投资银行的资深半导体行业分析师。请详细分析这份卖方研究报告。

请用中文输出结构化 Markdown，包含以下章节：

## 报告摘要
- 公司：[公司全称]
- 股票代码：[代码]
- 评级：[报告原始评级]
- 目标价：[新目标价] [货币] (此前：[旧目标价] [货币])
- 报告日期：[日期]

## 核心发现
- [发现一]
- [发现二]
- [发现三]

## 评级与目标价
- [评级理由及变动说明]
- 目标价：新旧对比、变动幅度、上行/下行空间
- 估值方法（P/E、EV/EBITDA、DCF、SOTP等）

## 营收与盈利预测
用表格呈现：
- FY2025 营收：[数值] [单位]，[同比增速%]
- FY2026 营收：[数值] [单位]，[同比增速%]
- FY2027 营收：[数值] [单位]，[同比增速%]
- EPS 预测
- 预测变动的主要驱动因素

## 业务分析
- 各业务板块表现
- 主要客户/终端市场
- 竞争格局
- 技术/产品周期

## 风险信号
- [具体风险一]
- [具体风险二]
- [具体风险三]

## 机会信号
- [具体机会一]
- [具体机会二]
- [具体机会三]

## 逻辑推导链
每个评级/目标价变动背后的因果逻辑。逐条列出：
- 驱动因素：[什么在驱动变化？]
- 方向：看多/看空/中性
- 置信度：高（多数据点支撑）/ 中（有数据但不确定）/ 低（推测性）
- 数据支撑：[指标名] = [具体数值]（来源：[报告原文/投行名称/公司口径]）
- 产业链影响：本公司 → 上游供应商 → 下游客户 → 竞争对手，分别影响什么
- 与前期变化：[相对于上次报告或市场预期的变动方向与幅度]
- 参考前期：[引用的前期报告日期或数据]

## 投资结论
- 看多逻辑
- 看空逻辑
- 关键催化剂
- 综合建议：[买入/持有/卖出]

请用具体数字，引用报告原文数据。区分预测值与实际值。每个逻辑推导链必须标明数据来源。"""


# ============ 核心类 ============

class PDFReportAnalyzer:
    """PDF 报告分析器 - 文本提取 + OCR fallback + LLM"""

    def __init__(self, api_key: str = None):
        self.api_key = api_key or API_KEY
        if not self.api_key:
            raise ValueError(
                "API key not set.\n"
                "  export ANTHROPIC_AUTH_TOKEN=<your-key>\n"
                "  or pass --api-key <your-key>"
            )

        self.client = anthropic.Anthropic(
            api_key=self.api_key,
            base_url=BASE_URL
        )
        self.model = MODEL

    # ---- File → Text (PDF + PPTX + XLSX + Image) ----

    @staticmethod
    def _xlsx_to_text(file_path: str, max_chars: int = None) -> str:
        """从 Excel 提取表格数据为文本"""
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        texts = []
        total_chars = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            texts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_values = [str(c) if c is not None else "" for c in row]
                row_text = " | ".join(row_values)
                if row_text.strip():
                    texts.append(row_text)
                    total_chars += len(row_text)
                    if max_chars and total_chars >= max_chars:
                        break
            if max_chars and total_chars >= max_chars:
                break
        wb.close()
        text = "\n".join(texts)
        if max_chars and len(text) > max_chars:
            text = text[:max_chars] + "\n\n[... truncated ...]"
        return text

    @staticmethod
    def _image_to_text(file_path: str, max_chars: int = None) -> str:
        """从图片 OCR 提取文本"""
        from PIL import Image
        img = Image.open(file_path)
        # 大图先缩放以加速 OCR
        if img.width > 3000 or img.height > 3000:
            ratio = min(3000 / img.width, 3000 / img.height)
            new_size = (int(img.width * ratio), int(img.height * ratio))
            img = img.resize(new_size, Image.LANCZOS)
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        if max_chars and len(text) > max_chars:
            text = text[:max_chars] + "\n\n[... truncated ...]"
        return text

    # ---- PPTX ----

    @staticmethod
    def _pptx_to_text(file_path: str, max_slides: int = None,
                      max_chars: int = None) -> str:
        """从 PowerPoint 提取文本"""
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = prs.slides[:max_slides] if max_slides else prs.slides

        texts = []
        total_chars = 0
        for i, slide in enumerate(slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    slide_texts.append("\n".join(rows))

            slide_content = "\n".join(slide_texts)
            if slide_content.strip():
                texts.append(f"--- Slide {i + 1} ---\n{slide_content}")
                total_chars += len(slide_content)
                if max_chars and total_chars >= max_chars:
                    break

        text = "\n\n".join(texts)
        if max_chars and len(text) > max_chars:
            text = text[:max_chars] + "\n\n[... truncated ...]"
        return text

    # ---- PDF → Text (dual strategy) ----

    @staticmethod
    def _is_garbled(text: str) -> bool:
        """判断文本是否是乱码（直接提取失败时触发 OCR）"""
        if len(text) < 100:
            return True
        # 乱码特征：高频非打印字符或连续无空格的大写字母串
        alpha = re.findall(r'[a-zA-Z]', text)
        if not alpha:
            return False
        upper = sum(1 for c in alpha if c.isupper())
        # 连续大写字母比例过高 → 可能是编码错误
        return upper / len(alpha) > 0.8 and len(text) < 5000

    @staticmethod
    def _ocr_page(page, dpi: int = OCR_DPI) -> str:
        """OCR 单个 PDF 页面"""
        pix = page.get_pixmap(dpi=dpi)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        return pytesseract.image_to_string(img, lang="chi_sim+eng")

    def file_to_text(self, file_path: str, max_pages: int = None,
                     max_chars: int = None, force_ocr: bool = False) -> tuple[str, bool]:
        """统一入口：PDF/PPTX/XLSX/Image → 文本提取
        Returns: (text, used_ocr)
        """
        ext = Path(file_path).suffix.lower()
        if ext == ".pptx":
            return self._pptx_to_text(file_path, max_slides=max_pages, max_chars=max_chars), False
        if ext in (".xlsx", ".xlsm", ".csv"):
            return self._xlsx_to_text(file_path, max_chars=max_chars), False
        if ext in (".jpg", ".jpeg", ".png", ".bmp", ".webp"):
            return self._image_to_text(file_path, max_chars=max_chars), True
        return self._pdf_to_text(file_path, max_pages, max_chars, force_ocr)

    def _pdf_to_text(self, pdf_path: str, max_pages: int = None,
                     max_chars: int = None, force_ocr: bool = False) -> tuple[str, bool]:
        """从 PDF 提取文本，先尝试直接提取，失败则 OCR
        Returns: (text, used_ocr)
        """
        doc = fitz.open(pdf_path)
        pages = doc[:max_pages] if max_pages else doc

        # Step 1: 尝试直接文本提取
        if not force_ocr:
            native_text = ""
            for page in pages:
                native_text += page.get_text("text") + "\n"

            if not self._is_garbled(native_text) and len(native_text) > 200:
                doc.close()
                text = native_text
                if max_chars and len(text) > max_chars:
                    text = text[:max_chars] + "\n\n[... truncated ...]"
                return text, False

        # Step 2: OCR fallback
        texts = []
        total_chars = 0
        for i, page in enumerate(pages):
            page_text = self._ocr_page(page)
            texts.append(f"--- Page {i + 1} ---\n{page_text}")
            total_chars += len(page_text)
            if max_chars and total_chars >= max_chars:
                break

        doc.close()
        text = "\n\n".join(texts)
        if max_chars and len(text) > max_chars:
            text = text[:max_chars] + "\n\n[... truncated ...]"
        return text, True

    # ---- API 调用 ----

    def _call_llm(self, system_prompt: str, report_text: str,
                  max_tokens: int = 4096) -> tuple[str, dict]:
        """调用 LLM (纯文本)"""
        # Sanitize: remove null bytes and non-printable control chars
        report_text = report_text.replace('\x00', '')
        report_text = ''.join(c for c in report_text if c.isprintable() or c in '\n\t\r')
        response = self.client.messages.create(
            model=self.model,
            system=system_prompt,
            messages=[{"role": "user", "content": report_text}],
            max_tokens=max_tokens,
            thinking={"type": "disabled"}
        )

        text_blocks = [b.text for b in response.content if b.type == "text"]
        text = "\n".join(text_blocks)
        usage = {
            "input_tokens": response.usage.input_tokens,
            "output_tokens": response.usage.output_tokens,
            "total_tokens": response.usage.input_tokens + response.usage.output_tokens
        }
        return text, usage

    # ---- Scanner ----

    def scan(self, pdf_path: str) -> dict:
        """快速扫描：前 2 页，判断相关性"""
        name = Path(pdf_path).name
        print(f"  🔍 Scanning: {name}")

        text, used_ocr = self.file_to_text(
            pdf_path, max_pages=2, max_chars=SCANNER_MAX_CHARS
        )
        method = "OCR" if used_ocr else "native"
        print(f"     Text: {len(text)} chars ({method})")

        user_prompt = (
            f'Analyze this report. The filename is "{name}". '
            f"Use the filename as a hint for the company name if the cover page is unclear.\n\n"
            f"--- REPORT TEXT ---\n{text}"
        )
        markdown, usage = self._call_llm(
            SYSTEM_PROMPT_SCANNER, user_prompt, max_tokens=1024
        )

        score_match = re.search(r'Relevance\s*Score[:\*\s]*(\d{1,2})', markdown)
        score = int(score_match.group(1)) if score_match else 0
        verdict = "DEEP_ANALYSIS" if score >= SCORE_THRESHOLD else "SKIP"

        print(f"     Score: {score}/20 → {verdict}  "
              f"(tokens: {usage['total_tokens']})")

        return {
            "score": score,
            "verdict": verdict,
            "markdown": markdown,
            "usage": usage,
            "ocr_used": used_ocr
        }

    # ---- Deep Analysis ----

    def deep_analyze(self, pdf_path: str, max_pages: int = None) -> dict:
        """深度分析：全文结构化输出"""
        name = Path(pdf_path).name

        text, used_ocr = self.file_to_text(
            pdf_path, max_pages=max_pages, max_chars=DEEP_MAX_CHARS
        )
        method = "OCR" if used_ocr else "native"
        print(f"  🔬 Deep analyzing: {name} ({len(text)} chars, {method})")

        markdown, usage = self._call_llm(
            SYSTEM_PROMPT_DEEP, text, max_tokens=4096
        )

        total = usage["total_tokens"]
        print(f"     Tokens: {total}  ~${total * 0.28 / 1_000_000:.4f}")

        return {"markdown": markdown, "usage": usage, "ocr_used": used_ocr}

    # ---- 完整流程 ----

    def analyze(self, pdf_path: str, auto_scan: bool = True,
                max_pages: int = None) -> dict:
        """完整分析：scan → deep (if relevant) → parse"""
        result = {
            "pdf_path": str(pdf_path),
            "pdf_name": Path(pdf_path).name,
            "analyzed_at": datetime.now().isoformat(),
            "mode": "auto"
        }

        if auto_scan:
            scan_result = self.scan(pdf_path)
            result["scan"] = scan_result

            if scan_result["verdict"] == "SKIP":
                result["mode"] = "scan_only"
                result["markdown"] = scan_result["markdown"]
                return result

        deep_result = self.deep_analyze(pdf_path, max_pages=max_pages)
        result["mode"] = "full" if auto_scan else "deep_only"
        result["markdown"] = deep_result["markdown"]
        result["usage"] = deep_result["usage"]

        try:
            from vision_parser import parse_vision_output
            parsed = parse_vision_output(result["markdown"])
            result["parsed"] = parsed
            result["alert_severity"] = parsed.get("alert_severity")
            result["alert_title"] = parsed.get("alert_title")
        except ImportError:
            pass

        return result

    def analyze_and_save(self, pdf_path: str, output_dir: str = None,
                         auto_scan: bool = True) -> dict:
        """分析并保存 Markdown + JSON"""
        result = self.analyze(pdf_path, auto_scan=auto_scan)

        output_dir = Path(output_dir or Path(pdf_path).parent)
        output_dir.mkdir(parents=True, exist_ok=True)
        stem = Path(pdf_path).stem

        md_path = output_dir / f"{stem}_analysis.md"
        md_path.write_text(result.get("markdown", ""), encoding="utf-8")
        print(f"     📄 {md_path}")

        json_path = output_dir / f"{stem}_analysis.json"
        serializable = {
            k: v for k, v in result.items()
            if k != "markdown" and not callable(v)
        }
        json_path.write_text(
            json.dumps(serializable, ensure_ascii=False, indent=2, default=str),
            encoding="utf-8", errors="replace"
        )
        print(f"     📋 {json_path}")

        return result


# ============ 批量处理 ============

def batch_analyze(directory: str, api_key: str = None,
                  output_dir: str = None, scan_only: bool = False,
                  no_scan: bool = False):
    """批量分析目录下所有 PDF"""
    analyzer = PDFReportAnalyzer(api_key=api_key)
    pdf_files = sorted(Path(directory).rglob("*.pdf"))
    pdf_files = [f for f in pdf_files
                 if not f.name.endswith("_analysis.md")]

    print(f"\n📊 Batch Analysis: {len(pdf_files)} PDFs in {directory}")
    print("=" * 60)

    ok = fail = skipped = 0
    for i, pdf_path in enumerate(pdf_files, 1):
        print(f"\n[{i}/{len(pdf_files)}] {pdf_path.name}")
        try:
            if scan_only:
                result = analyzer.scan(str(pdf_path))
                if result["verdict"] == "SKIP":
                    skipped += 1
                else:
                    ok += 1
            elif no_scan:
                analyzer.analyze_and_save(str(pdf_path), output_dir=output_dir)
                ok += 1
            else:
                result = analyzer.analyze_and_save(
                    str(pdf_path), output_dir=output_dir
                )
                if result.get("mode") == "scan_only":
                    skipped += 1
                else:
                    ok += 1
        except Exception as e:
            print(f"     ❌ Error: {e}")
            fail += 1

    print(f"\n{'=' * 60}")
    print(f"Done: {ok} ok, {skipped} skipped, {fail} failed")


# ============ CLI ============

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="PDF Report Analyzer - Native text + OCR fallback + LLM"
    )
    parser.add_argument("path", nargs="?", help="PDF file or directory path")
    parser.add_argument("--api-key", help="API key (or set ANTHROPIC_AUTH_TOKEN env)")
    parser.add_argument("--no-scan", action="store_true",
                        help="Skip scanner, deep analysis directly")
    parser.add_argument("--scan-only", action="store_true",
                        help="Scanner only, skip deep analysis")
    parser.add_argument("--max-pages", type=int, default=None,
                        help="Max pages for deep analysis")
    parser.add_argument("--batch", action="store_true",
                        help="Batch process all PDFs in directory")
    parser.add_argument("--output-dir", help="Output directory for results")
    parser.add_argument("--force-ocr", action="store_true",
                        help="Force OCR even if native text is available")

    args = parser.parse_args()

    if not args.path:
        parser.print_help()
        sys.exit(1)

    path = Path(args.path)

    if args.batch or path.is_dir():
        batch_analyze(
            str(path), api_key=args.api_key,
            output_dir=args.output_dir,
            scan_only=args.scan_only,
            no_scan=args.no_scan
        )
    else:
        if not path.exists():
            print(f"❌ File not found: {path}")
            sys.exit(1)

        analyzer = PDFReportAnalyzer(api_key=args.api_key)

        if args.scan_only:
            result = analyzer.scan(str(path))
            print(f"\n{result['markdown']}")
            # Save scan result
            output_dir = Path(args.output_dir or Path(path).parent)
            stem = Path(path).stem
            json_path = output_dir / f"{stem}_analysis.json"
            json_path.write_text(json.dumps(result, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
            md_path = output_dir / f"{stem}_analysis.md"
            md_path.write_text(result.get("markdown", ""), encoding="utf-8")
            print(f"     📄 {md_path}")
        elif args.no_scan:
            result = analyzer.analyze_and_save(
                str(path), output_dir=args.output_dir, auto_scan=False
            )
        else:
            result = analyzer.analyze_and_save(
                str(path), output_dir=args.output_dir
            )

        if result.get("parsed"):
            print()
            from vision_parser import print_summary
            print_summary(result["parsed"])
